#!/usr/bin/env python3
"""Export selected slide shape metadata from a PPTX into JSON.

Usage:
  python export_selected_slides.py source.pptx 13 70 71
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from pptx import Presentation

EMU = 914400


def to_in(value: int | float) -> float:
    return round(float(value) / EMU, 4)


def shape_record(shape):
    rec = {
        "shape_type": str(shape.shape_type),
        "left": to_in(shape.left),
        "top": to_in(shape.top),
        "width": to_in(shape.width),
        "height": to_in(shape.height),
        "name": getattr(shape, "name", ""),
    }
    if hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            rec["text"] = text
    return rec


def main(argv: list[str]) -> int:
    if len(argv) < 3:
        print(__doc__)
        return 1
    pptx_path = Path(argv[1])
    slide_numbers = [int(x) for x in argv[2:]]
    prs = Presentation(str(pptx_path))
    export = {}
    for n in slide_numbers:
        slide = prs.slides[n - 1]
        export[str(n)] = [shape_record(s) for s in slide.shapes]
    out_path = pptx_path.with_suffix(".selected-shapes.json")
    out_path.write_text(json.dumps(export, indent=2), encoding="utf-8")
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
